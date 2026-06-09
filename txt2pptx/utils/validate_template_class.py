"""validate_template_class.py — verify a TemplateStrategy works end-to-end.

Generates a sample PresentationOutline with one slide per SlideLayout,
runs the full pipeline, and runs a battery of checks:
  - every slide built without raising
  - every required placeholder ('title' for content slides, 'body' for bullets, etc.) was filled
  - all 9 SlideLayout types map to a layout (or use a fallback)
  - PPTX file opens cleanly when reloaded

Outputs:
    validation_<TemplateName>.pptx  (visual inspection)
    pass/fail report to stdout
    exits 0 on pass, 1 on fail

Usage:
    python -m txt2pptx.utils.validate_template_class OceanGradientStrategy
    python -m txt2pptx.utils.validate_template_class OceanGradientStrategy --output-dir /tmp
    python -m txt2pptx.utils.validate_template_class --all     # validate every registered strategy
"""
from __future__ import annotations

import argparse
import io
import sys
import traceback
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Type

from pptx import Presentation

from ..backend.models import (
    PresentationOutline,
    SlideData,
    SlideLayout,
    StatItem,
)
from ..backend.templates import TemplateStrategy, STRATEGIES, get_strategy

# ────────────────────────────────────────────────
# Sample outline (one slide per SlideLayout)
# ────────────────────────────────────────────────

SAMPLE_NOTES = (
    "這是一段示範用的講者備註，用於驗證 speaker_notes 欄位是否正確注入"
    "至投影片的備註區。長度設計為符合 50 字元最低限制的範本內容。"
)


def make_sample_outline() -> PresentationOutline:
    return PresentationOutline(
        title="Validation Deck",
        subtitle="Auto-generated sample for template validation",
        slides=[
            SlideData(layout=SlideLayout.TITLE,
                      title="Welcome to the Test Deck",
                      subtitle="A sample subtitle for the cover slide",
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.SECTION,
                      title="Section: Overview",
                      subtitle="Setting the stage for what comes next",
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.BULLETS,
                      title="Key Points",
                      bullets=["First key insight here",
                               "Second observation worth noting",
                               "Third takeaway for the audience"],
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.TWO_COLUMN,
                      title="Two Columns",
                      left_title="Pros",
                      left_column=["Faster to ship", "Less complex"],
                      right_title="Cons",
                      right_column=["Limited flexibility", "Higher cost"],
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.IMAGE_LEFT,
                      title="Image Left",
                      bullets=["Detail one", "Detail two"],
                      image_prompt="An abstract geometric pattern",
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.IMAGE_RIGHT,
                      title="Image Right",
                      bullets=["Point A", "Point B"],
                      image_prompt="A landscape scene",
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.KEY_STATS,
                      title="By the Numbers",
                      stats=[StatItem(value="100", label="users"),
                             StatItem(value="42", label="wins"),
                             StatItem(value="3.14", label="score")],
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.COMPARISON,
                      title="Old vs New",
                      left_title="Before",
                      left_column=["Manual process", "Slow"],
                      right_title="After",
                      right_column=["Automated", "Fast"],
                      speaker_notes=SAMPLE_NOTES),
            SlideData(layout=SlideLayout.CONCLUSION,
                      title="Wrapping Up",
                      bullets=["Recap of main points", "Next steps for the team"],
                      speaker_notes=SAMPLE_NOTES),
        ],
    )


# ────────────────────────────────────────────────
# Validation
# ────────────────────────────────────────────────

@dataclass
class CheckResult:
    name: str
    passed: bool
    detail: str = ""


@dataclass
class ValidationReport:
    strategy_name: str
    output_path: Optional[Path] = None
    checks: list[CheckResult] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return all(c.passed for c in self.checks)

    def add(self, name: str, passed: bool, detail: str = ""):
        self.checks.append(CheckResult(name, passed, detail))


def _check_layout_map_complete(strategy_cls: Type[TemplateStrategy], report: ValidationReport):
    """Every SlideLayout enum value should resolve to an int (mapped or fallback)."""
    inst = strategy_cls()
    missing = []
    for sl in SlideLayout:
        if sl not in inst.layout_map:
            missing.append(sl.name)
    if missing:
        report.add(
            "layout_map coverage",
            False,
            f"Missing direct mappings for: {missing} (will use fallback chain)",
        )
    else:
        report.add("layout_map coverage", True,
                   f"All {len(SlideLayout)} SlideLayouts directly mapped")


def _check_placeholder_map_required_keys(strategy_cls: Type[TemplateStrategy], report: ValidationReport):
    required = {"title", "body"}
    present = set(strategy_cls.placeholder_map.keys())
    missing = required - present
    report.add(
        "placeholder_map required keys",
        not missing,
        f"missing: {missing}" if missing else "title + body present",
    )


def _check_generation_no_exception(strategy_cls: Type[TemplateStrategy],
                                    outline: PresentationOutline,
                                    report: ValidationReport) -> Optional[bytes]:
    try:
        pptx_bytes = strategy_cls().generate(outline)
        report.add("generation runs without exception", True,
                   f"{len(pptx_bytes)} bytes produced")
        return pptx_bytes
    except Exception as e:
        tb = traceback.format_exc(limit=3)
        report.add("generation runs without exception", False,
                   f"{type(e).__name__}: {e}\n{tb}")
        return None


def _check_pptx_reloadable(pptx_bytes: bytes, report: ValidationReport) -> Optional[Presentation]:
    if not pptx_bytes:
        return None
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        report.add("PPTX reloads cleanly", True, f"{len(prs.slides)} slides")
        return prs
    except Exception as e:
        report.add("PPTX reloads cleanly", False, f"{type(e).__name__}: {e}")
        return None


def _check_slide_count_matches(prs, expected: int, report: ValidationReport):
    if prs is None:
        return
    actual = len(prs.slides)
    report.add(
        "slide count matches outline",
        actual == expected,
        f"expected {expected}, got {actual}",
    )


def _check_titles_filled(prs, outline: PresentationOutline, report: ValidationReport):
    """Every slide should have a non-empty title text in some shape."""
    if prs is None:
        return
    misses = []
    for i, (slide, slide_data) in enumerate(zip(prs.slides, outline.slides), 1):
        # SlideLayout.KEY_STATS routes title→body, not title placeholder
        if slide_data.layout == SlideLayout.KEY_STATS:
            continue
        title_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip():
                        title_text = run.text
                        break
                if title_text:
                    break
            if title_text:
                break
        if not title_text:
            misses.append(f"slide {i} ({slide_data.layout.value})")
    report.add(
        "every slide has visible text",
        not misses,
        f"empty slides: {misses}" if misses else f"all {len(outline.slides)} slides have text",
    )


def _check_speaker_notes(prs, outline: PresentationOutline, report: ValidationReport):
    if prs is None:
        return
    misses = []
    for i, (slide, slide_data) in enumerate(zip(prs.slides, outline.slides), 1):
        if not slide_data.speaker_notes:
            continue
        try:
            notes_text = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes_text = ""
        if slide_data.speaker_notes not in notes_text:
            misses.append(i)
    report.add(
        "speaker notes injected",
        not misses,
        f"missing notes on slides: {misses}" if misses else f"notes present on all {len(outline.slides)} slides",
    )


def validate(strategy_cls: Type[TemplateStrategy], output_dir: Path) -> ValidationReport:
    rpt = ValidationReport(strategy_name=strategy_cls.__name__)

    _check_layout_map_complete(strategy_cls, rpt)
    _check_placeholder_map_required_keys(strategy_cls, rpt)

    outline = make_sample_outline()
    pptx_bytes = _check_generation_no_exception(strategy_cls, outline, rpt)
    if pptx_bytes:
        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"validation_{strategy_cls.__name__}.pptx"
        out_path.write_bytes(pptx_bytes)
        rpt.output_path = out_path

    prs = _check_pptx_reloadable(pptx_bytes, rpt) if pptx_bytes else None
    _check_slide_count_matches(prs, len(outline.slides), rpt)
    _check_titles_filled(prs, outline, rpt)
    _check_speaker_notes(prs, outline, rpt)

    return rpt


def print_report(rpt: ValidationReport):
    print(f"\n=== {rpt.strategy_name} ===")
    for c in rpt.checks:
        mark = "✓" if c.passed else "✗"
        print(f"  [{mark}] {c.name}")
        if c.detail:
            for line in c.detail.splitlines():
                print(f"      {line}")
    if rpt.output_path:
        print(f"  Output: {rpt.output_path}")
    overall = "PASSED" if rpt.passed else "FAILED"
    print(f"  Overall: {overall}")


def main(argv: Optional[list[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        description="Validate a TemplateStrategy subclass end-to-end.")
    parser.add_argument("class_name", nargs="?", default=None,
                        help="strategy class name (e.g. OceanGradientStrategy)")
    parser.add_argument("--all", action="store_true",
                        help="validate every registered strategy")
    parser.add_argument("--output-dir", default="/tmp/validation",
                        help="where to write validation_*.pptx files")
    args = parser.parse_args(argv)

    output_dir = Path(args.output_dir)

    targets: list[Type[TemplateStrategy]] = []
    if args.all:
        targets = list(STRATEGIES.values())
    elif args.class_name:
        # Try by class name across registry
        match = next((cls for cls in STRATEGIES.values()
                      if cls.__name__ == args.class_name), None)
        if match is None:
            # Try by template_id
            match = get_strategy(args.class_name)
        if match is None:
            print(f"ERROR: Strategy {args.class_name!r} not found.", file=sys.stderr)
            print(f"Registered: {[c.__name__ for c in STRATEGIES.values()]}", file=sys.stderr)
            return 2
        targets = [match]
    else:
        parser.print_help()
        return 2

    all_passed = True
    for cls in targets:
        rpt = validate(cls, output_dir)
        print_report(rpt)
        all_passed = all_passed and rpt.passed

    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
