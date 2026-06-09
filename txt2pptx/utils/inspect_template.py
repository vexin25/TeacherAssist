"""inspect_template.py — read-only PPTX template introspection.

Usage:
    python -m txt2pptx.utils.inspect_template <template.pptx>           # human-readable
    python -m txt2pptx.utils.inspect_template <template.pptx> --json    # JSON

Outputs (human mode):
    - aspect ratio
    - slide masters & slide_layouts (name, idx, placeholders)
    - placeholder type, idx, dimensions, position
    - theme colors (primary/accent etc.)
    - default fonts

The same data is emitted as JSON when --json is passed, suitable
for downstream tooling (scaffold_template_class.py).
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, asdict, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Emu


def emu_to_inches(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    return round(Emu(emu).inches, 2)


@dataclass
class PlaceholderInfo:
    idx: int
    type: str
    name: str
    left_in: Optional[float] = None
    top_in: Optional[float] = None
    width_in: Optional[float] = None
    height_in: Optional[float] = None


@dataclass
class LayoutInfo:
    index: int
    name: str
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


@dataclass
class ThemeInfo:
    colors: dict[str, str] = field(default_factory=dict)
    major_font: Optional[str] = None
    minor_font: Optional[str] = None


@dataclass
class TemplateReport:
    file: str
    width_in: float
    height_in: float
    aspect_ratio: str
    n_masters: int
    n_layouts: int
    layouts: list[LayoutInfo] = field(default_factory=list)
    theme: ThemeInfo = field(default_factory=ThemeInfo)


# ────────────────────────────────────────────────
# Extraction
# ────────────────────────────────────────────────

def _placeholder_type_name(ph) -> str:
    pf = ph.placeholder_format
    try:
        return pf.type.name if pf.type is not None else "UNKNOWN"
    except Exception:
        return "UNKNOWN"


def _read_layout(layout, index: int) -> LayoutInfo:
    info = LayoutInfo(index=index, name=layout.name)
    for ph in layout.placeholders:
        try:
            info.placeholders.append(PlaceholderInfo(
                idx=ph.placeholder_format.idx,
                type=_placeholder_type_name(ph),
                name=ph.name,
                left_in=emu_to_inches(ph.left),
                top_in=emu_to_inches(ph.top),
                width_in=emu_to_inches(ph.width),
                height_in=emu_to_inches(ph.height),
            ))
        except Exception:
            # Some placeholder objects throw on attribute access; skip
            continue
    info.placeholders.sort(key=lambda p: p.idx)
    return info


def _read_theme(prs) -> ThemeInfo:
    """Best-effort theme extraction (colors + fonts)."""
    info = ThemeInfo()
    try:
        master = prs.slide_master
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from lxml import etree
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        tree = etree.fromstring(theme_part.blob)

        # Color scheme
        for clr in tree.findall(".//a:clrScheme/*", ns):
            tag = etree.QName(clr).localname  # e.g. "accent1"
            srgb = clr.find("./a:srgbClr", ns)
            sys_clr = clr.find("./a:sysClr", ns)
            if srgb is not None:
                info.colors[tag] = "#" + srgb.get("val")
            elif sys_clr is not None:
                info.colors[tag] = "#" + (sys_clr.get("lastClr") or sys_clr.get("val") or "")

        # Fonts
        major = tree.find(".//a:fontScheme/a:majorFont/a:latin", ns)
        minor = tree.find(".//a:fontScheme/a:minorFont/a:latin", ns)
        if major is not None:
            info.major_font = major.get("typeface")
        if minor is not None:
            info.minor_font = minor.get("typeface")
    except Exception:
        # Theme extraction is best-effort
        pass
    return info


def inspect(template_path: Path) -> TemplateReport:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))
    width_in = round(Emu(prs.slide_width).inches, 2)
    height_in = round(Emu(prs.slide_height).inches, 2)
    ratio = "16:9" if abs(width_in / height_in - 16 / 9) < 0.05 else (
            "4:3" if abs(width_in / height_in - 4 / 3) < 0.05 else f"{width_in}:{height_in}")

    report = TemplateReport(
        file=template_path.name,
        width_in=width_in,
        height_in=height_in,
        aspect_ratio=ratio,
        n_masters=len(prs.slide_masters),
        n_layouts=len(prs.slide_layouts),
    )
    for i, layout in enumerate(prs.slide_layouts):
        report.layouts.append(_read_layout(layout, i))
    report.theme = _read_theme(prs)
    return report


# ────────────────────────────────────────────────
# Reporters
# ────────────────────────────────────────────────

def report_human(r: TemplateReport) -> str:
    lines: list[str] = []
    lines.append(f"Template: {r.file}")
    lines.append(f"Slide size: {r.width_in} × {r.height_in} in  ({r.aspect_ratio})")
    lines.append(f"Slide masters: {r.n_masters}")
    lines.append(f"Layouts found: {r.n_layouts}")
    lines.append("")
    for lay in r.layouts:
        lines.append(f"  Layout {lay.index}: \"{lay.name}\"")
        if not lay.placeholders:
            lines.append("    (no placeholders)")
        for ph in lay.placeholders:
            dim = ""
            if ph.width_in is not None:
                dim = f", dim=({ph.width_in}, {ph.height_in}) at ({ph.left_in}, {ph.top_in})"
            lines.append(f"    - Placeholder idx={ph.idx}, type={ph.type}, name={ph.name!r}{dim}")
    lines.append("")
    if r.theme.colors:
        accent_colors = {k: v for k, v in r.theme.colors.items() if k.startswith("accent")}
        primary = r.theme.colors.get("dk2") or r.theme.colors.get("accent1")
        accent = r.theme.colors.get("accent2") or next(iter(accent_colors.values()), None)
        lines.append(f"Theme colors: primary={primary}, accent={accent}")
        lines.append(f"  All: {r.theme.colors}")
    if r.theme.major_font or r.theme.minor_font:
        lines.append(f"Default fonts: major={r.theme.major_font}, minor={r.theme.minor_font}")
    return "\n".join(lines)


def main(argv: Optional[list[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        description="Inspect a .pptx template and print its layouts/placeholders/theme.")
    parser.add_argument("template", help="path to .pptx file (or filename inside templates/)")
    parser.add_argument("--json", action="store_true", help="emit JSON instead of human text")
    args = parser.parse_args(argv)

    p = Path(args.template)
    if not p.exists():
        # Try templates/ dir
        templates_dir = Path(__file__).resolve().parent.parent / "templates"
        candidate = templates_dir / args.template
        if candidate.exists():
            p = candidate
        elif (templates_dir / f"{args.template}.pptx").exists():
            p = templates_dir / f"{args.template}.pptx"

    report = inspect(p)
    if args.json:
        print(json.dumps(asdict(report), indent=2, ensure_ascii=False))
    else:
        print(report_human(report))
    return 0


if __name__ == "__main__":
    sys.exit(main())
