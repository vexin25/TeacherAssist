"""TXT2PPTX FastAPI Application."""
import os
import uuid
import logging
from pathlib import Path
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware

from .models import GenerateRequest, GenerateResponse
from .llm_service import generate_outline
from .pptx_generator import generate_pptx as generate_pptx_code_drawn
from .pptx_generator_template import generate_pptx as generate_pptx_template_legacy
from .templates import get_strategy, is_registered, list_registered

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="TXT2PPTX", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
BASE_DIR = Path(__file__).parent.parent
GENERATED_DIR = BASE_DIR / "generated"
FRONTEND_DIR = BASE_DIR / "frontend"
GENERATED_DIR.mkdir(exist_ok=True)

# Serve frontend static files
app.mount("/static", StaticFiles(directory=str(FRONTEND_DIR)), name="static")


@app.get("/", response_class=HTMLResponse)
async def root():
    index_file = FRONTEND_DIR / "index.html"
    return index_file.read_text(encoding="utf-8")


@app.post("/api/generate", response_model=GenerateResponse)
async def generate_presentation(request: GenerateRequest):
    """Generate a PPTX presentation from text input."""
    try:
        # Step 1: Generate outline
        logger.info(f"Generating outline for {len(request.text)} chars, {request.num_slides} slides")
        outline = await generate_outline(request)
        logger.info(f"Outline generated: {outline.title}, {len(outline.slides)} slides")

        # Step 2: Generate PPTX
        # Dispatch order:
        #   1. code_drawn  → fully programmatic generator
        #   2. registered  → TemplateStrategy subclass (recommended path)
        #   3. legacy      → pptx_generator_template fallback (kept for backward compat)
        if request.template == "code_drawn":
            logger.info("Using code-drawn generator")
            pptx_bytes = generate_pptx_code_drawn(outline)
        elif is_registered(request.template):
            strategy_cls = get_strategy(request.template)
            logger.info(f"Using strategy: {strategy_cls.__name__}")
            pptx_bytes = strategy_cls().generate(outline)
        else:
            logger.warning(
                f"Template '{request.template}' is not a registered strategy; "
                f"falling back to legacy template loader. Registered: {list_registered()}"
            )
            pptx_bytes = generate_pptx_template_legacy(outline, template_id=request.template)

        # Step 3: Save file
        filename = f"{uuid.uuid4().hex[:8]}.pptx"
        filepath = GENERATED_DIR / filename
        filepath.write_bytes(pptx_bytes)
        logger.info(f"PPTX saved: {filepath} ({len(pptx_bytes)} bytes)")

        return GenerateResponse(
            success=True,
            filename=filename,
            message="簡報生成成功",
            outline=outline
        )

    except Exception as e:
        logger.error(f"Generation failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"生成失敗: {str(e)}")


@app.get("/api/download/{filename}")
async def download_file(filename: str):
    """Download generated PPTX file."""
    filepath = GENERATED_DIR / filename
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="檔案不存在")
    return FileResponse(
        path=str(filepath),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@app.get("/api/templates")
async def list_templates():
    """列出所有可用的簡報模板。"""
    # 模板中英文名稱對照表
    TEMPLATE_NAMES = {
        "College_Elegance": "學院典雅",
        "Data_Centric": "數據導向",
        "High_Contrast": "高調對比",
        "Minimalist_Corporate": "極簡商務",
        "Modernist": "摩登現代",
        "ocean_gradient": "預設版面",
        "Startup_Edge": "新創活力",
        "Zen_Serenity": "靜謐禪意",
    }

    templates = [
        {
            "id": "code_drawn",
            "name": "經典繪製",
            "description": "完全程式化繪製，靈活性高",
            "available": True,
            "is_template": False
        }
    ]

    # 檢查 templates 目錄下的所有 .pptx 檔案
    templates_dir = BASE_DIR / "templates"
    if templates_dir.exists():
        for template_file in templates_dir.glob("*.pptx"):
            template_id = template_file.stem
            try:
                # 嘗試載入模板驗證可用性
                from pptx import Presentation
                test_prs = Presentation(str(template_file))
                available = True
            except Exception as e:
                logger.warning(f"模板 {template_file.name} 不可用: {e}")
                available = False

            # 使用中文名稱對照表，若無對應則使用原始格式化名稱
            chinese_name = TEMPLATE_NAMES.get(template_id, template_id.replace("_", " ").title())

            templates.append({
                "id": template_id,
                "name": chinese_name,
                "description": f"使用 {chinese_name} 模板",
                "available": available,
                "is_template": True,
                "is_registered": is_registered(template_id),
            })

    return {"templates": templates}


# Onboarding message shown when an arbitrary .pptx is uploaded but no strategy
# class has been written for it yet. Mirrors the 5-step workflow documented in
# backend/templates/__init__.py.
ONBOARDING_INSTRUCTIONS_ZH = """\
⚠️ 自訂模板需要先註冊
您上傳的 "{filename}" 尚未註冊為支援的模板。

請執行以下步驟（預估 10–30 分鐘）：

1. 檢視模板結構（自動）：
   python -m txt2pptx.utils.inspect_template "{filename}"

2. 自動產生策略類別骨架：
   python -m txt2pptx.utils.scaffold_template_class "{filename}"

3. 編輯產生的 backend/templates/{snake}.py，
   解決所有 TODO 註解（手動，最費時的一步）

4. 驗證類別運作正確：
   python -m txt2pptx.utils.validate_template_class {camel}Strategy

5. 在 backend/templates/__init__.py 的 STRATEGIES 字典加入：
   "{stem}": {camel}Strategy
"""


@app.post("/api/upload-template")
async def upload_template(file: UploadFile = File(...)):
    """Accept a user-uploaded .pptx template.

    Behaviour:
      - Saves the file to templates/.
      - If a TemplateStrategy is already registered for this stem, returns
        success and the template becomes available immediately.
      - Otherwise, returns 202 with onboarding instructions describing the
        5-step workflow to register the template.
    """
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="僅支援 .pptx 檔案")

    templates_dir = BASE_DIR / "templates"
    templates_dir.mkdir(exist_ok=True)
    target = templates_dir / Path(file.filename).name

    content = await file.read()
    target.write_bytes(content)
    logger.info(f"Template uploaded: {target} ({len(content)} bytes)")

    stem = target.stem

    # Quick sanity check: can python-pptx open it?
    try:
        from pptx import Presentation
        Presentation(str(target))
    except Exception as e:
        target.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"無效的 PPTX 檔案：{e}")

    if is_registered(stem):
        return {
            "success": True,
            "registered": True,
            "template_id": stem,
            "message": f"模板 '{stem}' 已就緒，可立即使用。",
        }

    # Build snake/camel forms for the onboarding message
    import re
    snake = re.sub(r"([a-z0-9])([A-Z])", r"\1_\2", stem).lower()
    snake = re.sub(r"_+", "_", snake.replace("-", "_").replace(" ", "_")).strip("_")
    camel = "".join(p.capitalize() for p in re.split(r"[_\s\-]+", snake) if p)

    return JSONResponse(
        status_code=202,
        content={
            "success": True,
            "registered": False,
            "template_id": stem,
            "message": ONBOARDING_INSTRUCTIONS_ZH.format(
                filename=file.filename, snake=snake, camel=camel, stem=stem,
            ),
            "steps": [
                {"cmd": f'python -m txt2pptx.utils.inspect_template "{file.filename}"',
                 "desc": "檢視模板結構"},
                {"cmd": f'python -m txt2pptx.utils.scaffold_template_class "{file.filename}"',
                 "desc": "自動產生策略類別骨架"},
                {"cmd": f"# 編輯 backend/templates/{snake}.py，解決 TODO 註解",
                 "desc": "手動微調策略類別"},
                {"cmd": f"python -m txt2pptx.utils.validate_template_class {camel}Strategy",
                 "desc": "驗證類別運作"},
                {"cmd": f'# 在 backend/templates/__init__.py 加入 "{stem}": {camel}Strategy',
                 "desc": "註冊到 STRATEGIES"},
            ],
        },
    )


@app.get("/api/health")
async def health():
    return {"status": "ok", "version": "0.1.0"}
