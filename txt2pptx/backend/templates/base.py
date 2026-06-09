"""TemplateStrategy abstract base class.

Defines the contract every PPTX template strategy must implement.
Each concrete strategy encapsulates one template's recalculated parameters:
  - which slide_layout index maps to each SlideLayout
  - which placeholder idx is title/body/picture/etc.
  - font scale relative to the canonical ocean_gradient baseline
  - aspect ratio (used for warnings if mixed with other templates)

Subclasses override fill_* methods only when the default behaviour
(driven by placeholder_map + font_scale) is insufficient.
"""
from __future__ import annotations

import io
import logging
from abc import ABC
from pathlib import Path
from typing import ClassVar, Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from ..models import PresentationOutline, SlideData, SlideLayout

logger = logging.getLogger(__name__)

TEMPLATES_DIR = Path(__file__).resolve().parent.parent.parent / "templates"


class TemplateStrategy(ABC):
    """Base class for a PPTX template strategy.

    Concrete subclasses set the class-level configuration attributes
    and may override any fill_* method to customise per-template behaviour.
    """

    # ────────────────────────────────────────────────
    # Class-level configuration (subclasses override)
    # ────────────────────────────────────────────────

    #: filename inside templates/ (e.g. "ocean_gradient.pptx")
    template_file: ClassVar[str] = ""

    #: SlideLayout enum → slide_layouts index in the .pptx
    layout_map: ClassVar[dict[SlideLayout, int]] = {}

    #: Semantic name → placeholder idx within a layout.
    #: Required keys: 'title', 'body'. Optional: 'body_right', 'body_col2',
    #: 'body_col3', 'picture', 'slide_num'.
    placeholder_map: ClassVar[dict[str, int]] = {
        "title": 0,
        "body": 1,
        "body_right": 2,
        "body_col2": 3,
        "body_col3": 4,
        "picture": 10,
        "slide_num": 12,
    }

    #: Multiplier applied to all hard-coded font sizes (relative to ocean_gradient).
    #: e.g. font_scale=0.85 means a Pt(28) becomes Pt(23.8) for this template.
    font_scale: ClassVar[float] = 1.0

    #: Slide dimensions in inches (width, height). Used for warnings when
    #: a layout's content doesn't fit.
    aspect_ratio: ClassVar[tuple[float, float]] = (13.33, 7.5)

    # ────────────────────────────────────────────────
    # Public API
    # ────────────────────────────────────────────────

    def generate(self, outline: PresentationOutline) -> bytes:
        """Build a PPTX from the outline using this template strategy."""
        template_path = TEMPLATES_DIR / self.template_file
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        logger.info(f"[{self.__class__.__name__}] Loading {template_path.name}")
        prs = Presentation(str(template_path))
        self._clean_template_slides(prs)

        total = len(outline.slides)
        for idx, slide_data in enumerate(outline.slides, 1):
            dispatcher = self._builder_for(slide_data.layout)
            dispatcher(prs, slide_data, idx, total)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    # ────────────────────────────────────────────────
    # Layout dispatch
    # ────────────────────────────────────────────────

    def _builder_for(self, layout: SlideLayout):
        """Return the fill method for a given SlideLayout."""
        return {
            SlideLayout.TITLE: self.fill_title_slide,
            SlideLayout.SECTION: self.fill_section_header,
            SlideLayout.BULLETS: self.fill_bullets_slide,
            SlideLayout.TWO_COLUMN: self.fill_two_column,
            SlideLayout.IMAGE_LEFT: self.fill_image_left,
            SlideLayout.IMAGE_RIGHT: self.fill_image_right,
            SlideLayout.KEY_STATS: self.fill_key_stats,
            SlideLayout.COMPARISON: self.fill_comparison,
            SlideLayout.CONCLUSION: self.fill_conclusion,
        }.get(layout, self.fill_bullets_slide)

    def _layout_index(self, layout: SlideLayout) -> int:
        """Resolve a SlideLayout to the actual slide_layouts index, with fallback chain."""
        idx = self.layout_map.get(layout)
        if idx is not None:
            return idx
        # Fallback chain: TITLE→SECTION→BULLETS, IMAGE→BULLETS, COMPARISON→TWO_COLUMN, etc.
        fallback = {
            SlideLayout.TITLE: SlideLayout.SECTION,
            SlideLayout.SECTION: SlideLayout.BULLETS,
            SlideLayout.IMAGE_LEFT: SlideLayout.BULLETS,
            SlideLayout.IMAGE_RIGHT: SlideLayout.BULLETS,
            SlideLayout.KEY_STATS: SlideLayout.TWO_COLUMN,
            SlideLayout.COMPARISON: SlideLayout.TWO_COLUMN,
            SlideLayout.TWO_COLUMN: SlideLayout.BULLETS,
            SlideLayout.CONCLUSION: SlideLayout.BULLETS,
        }.get(layout)
        if fallback is not None:
            logger.debug(f"Layout {layout} not mapped; falling back to {fallback}")
            return self._layout_index(fallback)
        # Last resort: first layout
        return 0

    # ────────────────────────────────────────────────
    # Default fill_* implementations
    # ────────────────────────────────────────────────

    def fill_title_slide(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.TITLE)
        self._safe_fill(slide, "title", slide_data.title)
        self._safe_fill(slide, "body", slide_data.subtitle)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_section_header(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.SECTION)
        self._safe_fill(slide, "title", slide_data.title)
        self._safe_fill(slide, "body", slide_data.subtitle)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_bullets_slide(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.BULLETS)
        self._safe_fill(slide, "title", slide_data.title)
        self._fill_bullets(slide, "body", slide_data.bullets)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_two_column(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.TWO_COLUMN)
        self._safe_fill(slide, "title", slide_data.title)
        self._fill_column(slide, "body", slide_data.left_title, slide_data.left_column)
        self._fill_column(slide, "body_right", slide_data.right_title, slide_data.right_column)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_image_left(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.IMAGE_LEFT)
        self._safe_fill(slide, "title", slide_data.title)
        self._fill_bullets(slide, "body", slide_data.bullets)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_image_right(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.IMAGE_RIGHT)
        self._safe_fill(slide, "title", slide_data.title)
        self._fill_bullets(slide, "body", slide_data.bullets)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_key_stats(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.KEY_STATS)
        # Some templates use 'body' for the section title in stats layouts
        self._safe_fill(slide, "body", slide_data.title)
        stats = (slide_data.stats or [])[:3]
        for i, stat in enumerate(stats):
            ph_name = ["body_right", "body_col2", "body_col3"][i]
            self._format_stat(slide, ph_name, stat)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_comparison(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.COMPARISON)
        self._safe_fill(slide, "title", slide_data.title)
        self._fill_column(slide, "body", slide_data.left_title, slide_data.left_column)
        self._fill_column(slide, "body_right", slide_data.right_title, slide_data.right_column)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    def fill_conclusion(self, prs, slide_data: SlideData, idx: int, total: int):
        slide = self._add_slide(prs, SlideLayout.CONCLUSION)
        self._safe_fill(slide, "title", slide_data.title)
        if slide_data.bullets:
            self._fill_bullets(slide, "body", slide_data.bullets)
        elif slide_data.subtitle:
            self._safe_fill(slide, "body", slide_data.subtitle)
        self._fill_slide_number(slide, idx, total)
        self._fill_speaker_notes(slide, slide_data.speaker_notes)

    # ────────────────────────────────────────────────
    # Low-level helpers (shared by all subclasses)
    # ────────────────────────────────────────────────

    def _add_slide(self, prs, layout: SlideLayout):
        layout_idx = self._layout_index(layout)
        layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
        return prs.slides.add_slide(prs.slide_layouts[layout_idx])

    def _ph(self, slide, ph_name: str):
        """Resolve semantic ph_name → placeholder object, or None."""
        ph_idx = self.placeholder_map.get(ph_name)
        if ph_idx is None:
            return None
        try:
            return slide.placeholders[ph_idx]
        except KeyError:
            return None

    def _safe_fill(self, slide, ph_name: str, text: Optional[str]):
        if not text:
            return
        ph = self._ph(slide, ph_name)
        if ph is not None:
            ph.text = text

    def _fill_bullets(self, slide, ph_name: str, items: Optional[list[str]]):
        ph = self._ph(slide, ph_name)
        if ph is None or not items:
            return
        tf = ph.text_frame
        tf.clear()
        for i, item in enumerate(items):
            if i == 0:
                tf.paragraphs[0].text = item
            else:
                p = tf.add_paragraph()
                p.text = item

    def _fill_column(self, slide, ph_name: str, title: Optional[str], items: Optional[list[str]]):
        ph = self._ph(slide, ph_name)
        if ph is None:
            return
        tf = ph.text_frame
        tf.clear()
        has_title = bool(title)
        if has_title:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.bold = True
        if items:
            for i, item in enumerate(items):
                if i == 0 and not has_title:
                    tf.paragraphs[0].text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item

    def _format_stat(self, slide, ph_name: str, stat):
        ph = self._ph(slide, ph_name)
        if ph is None or stat is None:
            return
        tf = ph.text_frame
        tf.clear()

        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        run_val = p_val.add_run()
        run_val.text = stat.value
        run_val.font.bold = True
        run_val.font.size = Pt(28 * self.font_scale)

        p_label = tf.add_paragraph()
        p_label.alignment = PP_ALIGN.CENTER
        run_label = p_label.add_run()
        run_label.text = stat.label
        run_label.font.size = Pt(11 * self.font_scale)

    def _fill_slide_number(self, slide, num: int, total: int):
        self._safe_fill(slide, "slide_num", f"{num} / {total}")

    def _fill_speaker_notes(self, slide, notes: Optional[str]):
        if not notes:
            return
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    @staticmethod
    def _clean_template_slides(prs):
        """Remove all pre-existing slides from the template."""
        xml_slides = prs.slides._sldIdLst
        for sldId in list(xml_slides):
            rId = sldId.get(qn("r:id"))
            prs.part.drop_rel(rId)
            xml_slides.remove(sldId)
